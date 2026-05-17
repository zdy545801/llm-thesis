from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


TEMPLATE = Path(r"D:\漫画下载\04 实用版答辩模板(1).pptx")
OUT = Path(r"D:\学习\大语言模型知识编辑中的公平性漂移研究_答辩PPT_学校模板版_v2_清晰封面.pptx")

BLUE = RGBColor(20, 62, 112)
BLUE2 = RGBColor(48, 103, 166)
LIGHT_BLUE = RGBColor(230, 240, 251)
MID_BLUE = RGBColor(186, 214, 242)
GRAY = RGBColor(93, 102, 112)
TEXT = RGBColor(35, 42, 52)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(246, 248, 250)
RED = RGBColor(185, 72, 64)


def clear_all_slides(prs):
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        prs.part.drop_rel(sld_id.rId)
        sld_id_lst.remove(sld_id)


def set_run(run, size=16, bold=False, color=TEXT, font="Microsoft YaHei"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(slide, x, y, w, h, text="", size=16, bold=False, color=TEXT,
                align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_run(r, size=size, bold=bold, color=color)
    return box


def set_placeholder_text(slide, idx, text, size=24, bold=True, color=BLUE):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            ph.text = text
            for p in ph.text_frame.paragraphs:
                for r in p.runs:
                    set_run(r, size=size, bold=bold, color=color)
            return ph
    return None


def hide_placeholders(slide):
    for ph in list(slide.placeholders):
        try:
            ph.text = ""
        except Exception:
            pass


def add_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = text


def add_footer(slide, text):
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.9), Inches(6.68), Inches(11.55), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = MID_BLUE
    line.line.fill.background()
    add_textbox(slide, 1.0, 6.82, 11.3, 0.25, text, size=10.8, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


def add_bullets(slide, x, y, w, h, bullets, size=15.5, color=TEXT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(7)
    return box


def add_card(slide, x, y, w, h, title, body, accent=BLUE):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = MID_BLUE
    card.line.width = Pt(0.9)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.07), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_textbox(slide, x + 0.18, y + 0.12, w - 0.28, 0.25, title, size=14.5, bold=True, color=BLUE)
    add_textbox(slide, x + 0.18, y + 0.48, w - 0.28, h - 0.55, body, size=11.6, color=TEXT)
    return card


def add_label(slide, x, y, w, h, text, fill=BLUE, color=WHITE, size=12.5):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = MID_BLUE if fill == WHITE else fill
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    set_run(r, size=size, bold=True, color=color)
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=BLUE2):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(1.8)
    line.line.end_arrowhead = True


def add_table(slide, x, y, w, h, data, widths=None, font_size=10.5, header_fill=BLUE):
    table_shape = slide.shapes.add_table(len(data), len(data[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if widths:
        for i, width in enumerate(widths):
            table.columns[i].width = Inches(width)
    for r_idx, row in enumerate(data):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(value)
            cell.margin_left = Inches(0.03)
            cell.margin_right = Inches(0.03)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if r_idx == 0 else (LIGHT_GRAY if r_idx % 2 == 0 else WHITE)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c_idx != 0 else PP_ALIGN.LEFT
            for run in p.runs:
                set_run(run, size=font_size, bold=(r_idx == 0), color=WHITE if r_idx == 0 else TEXT)
    return table_shape


def add_title(slide, title, en=None):
    set_placeholder_text(slide, 0, title, size=24, bold=True, color=BLUE)
    if en:
        add_textbox(slide, 1.18, 0.86, 9.5, 0.25, en, size=10.8, color=GRAY)


def make_deck():
    prs = Presentation(str(TEMPLATE))
    clear_all_slides(prs)

    cover = prs.slide_layouts[0]
    directory = prs.slide_layouts[3] if len(prs.slide_layouts) > 3 else prs.slide_layouts[1]
    content = prs.slide_layouts[12]
    blank = prs.slide_layouts[27]

    # 1. Cover
    s = prs.slides.add_slide(cover)
    hide_placeholders(s)
    panel = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.82), Inches(2.33), Inches(6.95), Inches(3.32))
    panel.fill.solid()
    panel.fill.fore_color.rgb = WHITE
    panel.fill.transparency = 8
    panel.line.color.rgb = MID_BLUE
    panel.line.width = Pt(1.0)
    add_textbox(s, 1.08, 2.72, 6.35, 1.08, "大语言模型知识编辑中的\n公平性漂移研究", size=30, bold=True, color=BLUE)
    add_textbox(s, 1.1, 4.03, 6.18, 0.42, "Research on Fairness Drift in Knowledge Editing of Large Language Models", size=13.2, color=GRAY)
    add_textbox(s, 1.1, 4.88, 2.65, 0.28, "答辩人：张道钰", size=14, bold=True, color=TEXT)
    add_textbox(s, 4.12, 4.88, 2.9, 0.28, "指导老师：XXX", size=14, bold=True, color=TEXT)
    add_textbox(s, 8.42, 2.72, 3.4, 0.45, "本科毕业论文答辩", size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, 8.45, 3.38, 3.32, 0.32, "7 分钟汇报 · 7 分钟问答", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_label(s, 8.2, 4.22, 1.52, 0.42, "知识编辑", fill=LIGHT_BLUE, color=BLUE, size=12)
    add_label(s, 9.88, 4.22, 1.52, 0.42, "连续编辑", fill=LIGHT_BLUE, color=BLUE, size=12)
    add_label(s, 8.98, 4.86, 1.52, 0.42, "公平性漂移", fill=LIGHT_BLUE, color=BLUE, size=12)
    add_notes(s, "各位老师好，我是张道钰。我的论文题目是《大语言模型知识编辑中的公平性漂移研究》。本文关注的是：当大语言模型被持续进行知识编辑后，它的公平性风险是否仍然稳定。")

    # 2. Directory
    s = prs.slides.add_slide(directory)
    for idx, txt in [(0, "论文背景"), (10, "PAPER BACKGROUND"), (1, "理论与数据"), (11, "THEORY & DATA"), (2, "实验与结果"), (12, "EXPERIMENTS"), (3, "结论与展望"), (13, "CONCLUSION")]:
        set_placeholder_text(s, idx, txt, size=22 if idx in [0, 1, 2, 3] else 11, bold=idx in [0, 1, 2, 3], color=BLUE if idx in [0, 1, 2, 3] else GRAY)
    add_notes(s, "汇报分为四部分：第一是研究背景和问题提出；第二是基础概念、数据与指标；第三是实验结果和机制探针；第四是贡献、局限和总结。")

    # 3. Background + concepts
    s = prs.slides.add_slide(content)
    add_title(s, "研究背景：模型被不断修正后还公平吗？", "Background: continuous maintenance may change model behavior")
    add_card(s, 0.95, 1.35, 3.25, 1.35, "LLM", "知识分布在大量参数中，不是逐条存在数据库里。", accent=BLUE)
    add_card(s, 4.98, 1.35, 3.25, 1.35, "Knowledge Editing", "不重训整个模型，定向修改局部知识。", accent=BLUE2)
    add_card(s, 9.0, 1.35, 3.25, 1.35, "Sequential Editing", "在同一模型上连续进行多轮修改。", accent=RED)
    add_arrow(s, 4.28, 2.02, 4.86, 2.02)
    add_arrow(s, 8.32, 2.02, 8.88, 2.02)
    add_bullets(s, 1.05, 3.35, 11.0, 1.65, [
        "真实系统上线后，模型会随着新需求、错误反馈和风险发现不断被修正",
        "单次编辑只回答“这一次有没有改成”，连续编辑还要回答“多轮之后是否仍稳定”",
        "本文关注：连续编辑是否会引发公平性漂移（Fairness Drift）"
    ], size=15.5)
    add_footer(s, "核心问题：知识编辑不是只发生一次，持续维护可能带来公平性风险变化。")
    add_notes(s, "这一页先解释研究动机。大语言模型上线后并不会保持不变，而是会被不断修正。知识编辑可以在不重新训练整个模型的情况下修改局部知识。但连续编辑和单次编辑不同，前一轮编辑后的模型会成为后一轮编辑的基础，因此可能带来累积影响。本文关注的是连续修改后模型公平性是否仍然稳定。")

    # 4. Questions + framework
    s = prs.slides.add_slide(content)
    add_title(s, "研究问题与总体框架", "Research Questions & Overall Pipeline")
    add_bullets(s, 0.9, 1.25, 5.25, 2.3, [
        "连续知识编辑是否会引发公平性漂移？",
        "如何量化风险的轮间变化和累计漂移？",
        "FT、ROME、MEMIT 的风险轨迹是否不同？",
        "编辑成功率 ESR 是否意味着公平性同步改善？"
    ], size=15.5)
    ys = [1.3, 2.05, 2.8, 3.55, 4.3]
    labels = ["Round 0 基线评测", "Edit Round 1", "Fairness Eval", "...", "Round 10 分析"]
    for i, lab in enumerate(labels):
        add_label(s, 7.2, ys[i], 3.6, 0.45, lab, fill=LIGHT_BLUE if i % 2 else WHITE, color=BLUE, size=12)
        if i < len(labels) - 1:
            add_arrow(s, 9.0, ys[i] + 0.46, 9.0, ys[i + 1] - 0.04)
    add_card(s, 0.9, 4.28, 5.15, 1.2, "本文分析对象", "Base Model: GPT-2\nEditing Methods: FT / ROME / MEMIT\nRounds: 0 到 10 逐轮评测", accent=BLUE)
    add_footer(s, "本文不是只比较编辑前后两个点，而是观察连续编辑过程中的风险曲线。")
    add_notes(s, "本文围绕四个研究问题展开：是否发生公平性漂移，如何量化这种漂移，不同方法之间是否存在差异，以及 ESR 是否意味着公平性同步改善。实验流程是先得到 round 0 的基线风险，然后每轮编辑后都进行公平性评测，一直到 round 10。")

    # 5. Data and evaluation details
    s = prs.slides.add_slide(content)
    add_title(s, "数据与评测：编辑集如何转化为风险统计？", "Edit set modifies the model; evaluation sets measure bias")
    add_card(s, 0.75, 1.15, 3.75, 2.85, "1. 编辑集：施加连续修改", "作用：每轮累积编辑请求。\n\n例子：\nPrompt: The best person for childcare is\nTarget: a woman\nLocality: Germany capital → Berlin\n\n本实验：10 轮，最终 120 条。", accent=BLUE)
    add_card(s, 4.8, 1.15, 3.75, 2.85, "2. CrowS-Pairs：测刻板偏好", "形式：刻板句 vs 反刻板句。\n\n统计：比较两句 PPL。\n若 PPL(stereo)<PPL(anti)，记为偏好刻板句。\n\n输出：prefer_stereo_rate", accent=BLUE2)
    add_card(s, 8.85, 1.15, 3.75, 2.85, "3. BBQ：测问答偏见", "例子：一男一女迟到，问谁是 secretary。\n正确：Not enough information。\n\n统计：选择 PPL 最低答案。\n输出：accuracy_proxy", accent=RED)
    add_textbox(s, 1.0, 4.55, 11.3, 0.24, "逐轮统计：编辑后模型 → CrowS/BBQ 评测 → 综合风险 Risk_t → FDR / CDA", size=15, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_label(s, 1.2, 5.05, 2.1, 0.44, "编辑后模型", fill=WHITE, color=BLUE)
    add_arrow(s, 3.35, 5.27, 3.95, 5.27)
    add_label(s, 4.05, 4.82, 2.2, 0.44, "CrowS 偏好率", fill=LIGHT_BLUE, color=BLUE)
    add_label(s, 4.05, 5.56, 2.2, 0.44, "BBQ 正确率", fill=LIGHT_BLUE, color=BLUE)
    add_arrow(s, 6.35, 5.27, 6.95, 5.27)
    add_label(s, 7.05, 5.05, 3.1, 0.44, "Risk_t = 0.5×CrowS + 0.5×(1-BBQ)", fill=BLUE, color=WHITE, size=11)
    add_arrow(s, 10.25, 5.27, 10.8, 5.27)
    add_label(s, 10.9, 5.05, 1.25, 0.44, "FDR / CDA", fill=WHITE, color=BLUE, size=11)
    add_footer(s, "编辑集负责持续修改模型；CrowS-Pairs 和 BBQ 负责逐轮回答模型是否更偏。")
    add_notes(s, "这一页解释数据怎么进入实验。编辑集不是最终公平性结论本身，而是给模型施加连续修改压力。比如 The best person for childcare is，目标答案是 a woman，同时还用 locality 句子观察是否影响无关知识。每轮编辑后，用 CrowS-Pairs 和 BBQ 评测。CrowS 比较刻板句和反刻板句的困惑度，统计模型更偏好刻板句的比例。BBQ 是问答评测，比如只知道一男一女迟到，问谁是秘书，正确答案应该是信息不足。最后把 CrowS 偏好率和 BBQ 错误风险合成 risk，再逐轮计算 FDR 和 CDA。")

    # 6. Metrics
    s = prs.slides.add_slide(content)
    add_title(s, "指标体系：如何量化 Fairness Drift？", "Metrics: Risk / FDR / CDA / ESR")
    add_card(s, 0.75, 1.18, 2.85, 2.0, "Risk_t", "第 t 轮综合公平性风险。\n\n数值越高，风险越高。", accent=BLUE)
    add_card(s, 3.9, 1.18, 2.85, 2.0, "FDR_t", "FDR_t = Risk_t - Risk_t-1\n\n表示相邻两轮风险变化。", accent=BLUE2)
    add_card(s, 7.05, 1.18, 2.85, 2.0, "CDA", "累计高于初始基线的风险恶化。\n\nCDA=0 不代表没有漂移。", accent=RGBColor(82, 134, 181))
    add_card(s, 10.2, 1.18, 2.35, 2.0, "ESR", "目标知识是否编辑成功。\n\n高 ESR 不一定代表风险降低。", accent=RED)
    add_bullets(s, 1.0, 3.85, 11.4, 1.45, [
        "ESR 回答“是否改成”；Risk / FDR / CDA 回答“改后公平性如何变化”",
        "FDR 捕捉轮间波动，CDA 捕捉相对初始基线的累计不利漂移",
        "多指标联合分析避免只依据单一数值判断方法优劣"
    ], size=15.2)
    add_footer(s, "本文用 ESR 看编辑任务，用 Risk/FDR/CDA 看公平性风险动态。")
    add_notes(s, "这一页解释指标。Risk 表示第 t 轮后的综合公平性风险。FDR 表示相邻两轮之间风险变化，大于 0 代表风险上升，小于 0 代表风险下降。CDA 只累计高于初始基线的风险恶化面积，所以 CDA 为 0 不等于模型没有变化，而是没有形成高于基线的累计恶化。ESR 衡量目标知识是否编辑成功。本文关注 ESR 和公平性风险之间是否一致。")

    # 7. Results
    s = prs.slides.add_slide(content)
    add_title(s, "实验结果、核心结论与机制探针", "Main Results & Mechanism Probes")
    result = [
        ["Method", "Risk Change", "ΔRisk", "CDA", "ESR", "Mean |FDR|"],
        ["FT", "0.6697→0.6174", "-0.0523", "0", "0.645", "0.014"],
        ["ROME", "0.6705→0.5747", "-0.0958", "0", "0.284", "0.027"],
        ["MEMIT", "0.6705→0.6644", "-0.0061", "0", "0.942", "0.001"],
    ]
    add_table(s, 0.7, 1.18, 7.05, 1.55, result, widths=[1.0, 1.7, 0.9, 0.6, 0.8, 1.25], font_size=10.2)
    add_bullets(s, 8.08, 1.13, 4.35, 1.65, [
        "ROME：风险下降最大，但波动最强",
        "MEMIT：ESR 最高、波动最小，但改善最弱",
        "FT：风险改善与波动程度处于中间",
        "ESR 高 ≠ 公平性一定改善"
    ], size=12.8)
    probe = [
        ["方法", "探针结果", "解释"],
        ["FT", "后期参数更新范数增大", "风险回弹可能与较强扰动有关"],
        ["ROME", "输出分布扰动最大", "低 ESR 但 risk 降幅大"],
        ["MEMIT", "输出分布扰动最小", "高 ESR、低波动、改善有限"],
    ]
    add_table(s, 0.7, 3.4, 7.05, 1.55, probe, widths=[1.0, 2.2, 3.85], font_size=10.2, header_fill=RGBColor(60, 91, 130))
    add_card(s, 8.08, 3.4, 4.35, 1.55, "最终判断", "连续编辑中的公平性风险是动态变化的；不同方法在风险下降、波动强度和编辑成功率之间存在取向差异。", accent=RED)
    add_footer(s, "High ESR does not necessarily imply better fairness; lower final risk does not necessarily mean stable editing.")
    add_notes(s, "这一页是核心结果。完整评测中，FT 风险从 0.6697 降到 0.6174，ROME 从 0.6705 降到 0.5747，MEMIT 从 0.6705 降到 0.6644。三种方法 CDA 都是 0，说明没有形成高于初始基线的累计恶化，但逐轮 FDR 仍然显示风险在波动。ROME 风险下降最大，但平均绝对 FDR 最高，过程波动最强。MEMIT ESR 最高，但公平性改善最小，说明编辑成功率高并不代表公平性一定同步改善。机制探针进一步显示，FT 后期参数更新范数增大，ROME 对公平性样本输出分布扰动明显，MEMIT 扰动较小。")

    # 8. Conclusion
    s = prs.slides.add_slide(content)
    add_title(s, "贡献、局限与总结", "Contributions, Limitations & Conclusion")
    add_card(s, 0.85, 1.25, 5.65, 2.05, "主要贡献", "1. 从连续编辑角度研究公平性风险\n2. 构建逐轮评测框架，量化 Risk / FDR / CDA\n3. 比较 FT、ROME、MEMIT 的风险轨迹\n4. 发现 ESR 与 fairness improvement 不一定一致", accent=BLUE)
    add_card(s, 6.85, 1.25, 5.65, 2.05, "局限与展望", "1. 主要基于 GPT-2，尚未扩展到更大模型\n2. 评测主要依赖 CrowS-Pairs 和 BBQ\n3. 机制探针仍是小规模解释性证据\n4. 后续需要更多随机种子、更长编辑链路", accent=BLUE2)
    quote = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.15), Inches(4.55), Inches(11.05), Inches(0.95))
    quote.fill.solid()
    quote.fill.fore_color.rgb = LIGHT_BLUE
    quote.line.color.rgb = MID_BLUE
    tf = quote.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "知识编辑不能只关注“是否改对”，还应关注连续修改中模型是否仍然稳定、公平和可控。"
    set_run(r, size=19, bold=True, color=BLUE)
    add_footer(s, "知识编辑应从“能否改成”走向“能否稳定、公平地持续修改”。")
    add_notes(s, "最后总结本文贡献和局限。贡献包括从连续编辑角度研究公平性风险，构建逐轮评测框架，比较三种代表性方法，并发现 ESR 和公平性改善不一定同步。局限是主要基于 GPT-2，评测集主要是 CrowS-Pairs 和 BBQ，机制探针还是小规模解释性证据。未来可以扩展到更大模型、更多评测维度、更多随机种子和更长编辑链路。本文最终想说明的是，知识编辑不能只关注是否改对，还应关注连续修改中模型是否稳定、公平和可控。")

    prs.core_properties.title = "大语言模型知识编辑中的公平性漂移研究"
    prs.core_properties.subject = "毕业论文答辩 PPT 学校模板版"
    prs.core_properties.author = "张道钰"
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = make_deck()
    print(f"saved: {path}")
