from pathlib import Path

from pptx import Presentation
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt


PPT = Path(r"D:\学习\单页_公平性指标定义.pptx")


def first_style(shape):
    if not hasattr(shape, "text_frame"):
        return {}
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if r.text:
                f = r.font
                return {
                    "name": f.name,
                    "size": f.size,
                    "bold": f.bold,
                    "italic": f.italic,
                    "color": f.color.rgb if f.color.type == 1 else None,
                }
    return {}


def clear_bullets(p):
    pPr = p._p.get_or_add_pPr()
    for child in list(pPr):
        if child.tag.endswith(("}buChar", "}buAutoNum", "}buBlip", "}buNone")):
            pPr.remove(child)
    pPr.insert(0, OxmlElement("a:buNone"))


def set_text(shape, text, size=None, bold=None):
    style = first_style(shape)
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        clear_bullets(p)
        p.text = ""
        run = p.add_run()
        run.text = line
        f = run.font
        if style.get("name"):
            f.name = style["name"]
        f.size = Pt(size) if size else style.get("size")
        f.bold = bold if bold is not None else style.get("bold")
        f.italic = style.get("italic")
        if style.get("color") is not None:
            f.color.rgb = style["color"]


prs = Presentation(PPT)
s = prs.slides[0]

set_text(s.shapes[0], "本文公平性指标定义：如何量化公平性漂移？", 25, True)
set_text(s.shapes[1], "Metric definitions: Risk, FDR, CDA and ESR", 13)

set_text(s.shapes[4], "1. Fairness Risk", 16, True)
set_text(
    s.shapes[5],
    "回答：第 t 轮后模型整体公平性风险多高？\n"
    "Risk_t = 0.5 x CrowS_t + 0.5 x (1 - BBQ_t)\n"
    "数值越高，说明刻板偏好或问答偏见风险越高。",
    12,
)

set_text(s.shapes[8], "2. FDR", 16, True)
set_text(
    s.shapes[9],
    "Fairness Drift Rate\n"
    "回答：相邻两轮之间风险怎么变？\n"
    "FDR_t = Risk_t - Risk_{t-1}\n"
    "FDR>0 表示风险上升；FDR<0 表示风险下降。",
    12,
)

set_text(s.shapes[12], "3. CDA", 16, True)
set_text(
    s.shapes[13],
    "Cumulative Drift Area\n"
    "回答：是否形成高于初始基线的累计恶化？\n"
    "CDA = sum max(0, Risk_t - Risk_0)\n"
    "CDA=0 不等于没有波动，只表示没有累计超过基线。",
    12,
)

set_text(
    s.shapes[16],
    "补充指标：ESR（Edit Success Rate）\n"
    "ESR = mean(post.rewrite_acc)，回答“目标知识有没有改成功”。\n"
    "但 ESR 衡量的是编辑任务成功率，不直接衡量公平性。\n\n"
    "因此本文同时看：ESR 是否改成功 + Risk/FDR/CDA 公平性是否稳定。",
    22,
)

set_text(s.shapes[18], "核心解释：ESR 看“改没改成”，Risk/FDR/CDA 看“连续修改后是否仍公平、稳定”。", 14)

prs.save(PPT)
print(PPT)
