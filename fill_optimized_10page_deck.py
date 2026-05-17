from pathlib import Path

from pptx import Presentation
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt


PPT = Path(r"D:\学习\222_讲稿优化10页版.pptx")
OUT = Path(r"D:\学习\222_讲稿优化10页完成版.pptx")


def first_style(shape):
    if not hasattr(shape, "text_frame"):
        return {}
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if r.text:
                f = r.font
                return {
                    "name": f.name,
                    "size": f.size,
                    "bold": f.bold,
                    "italic": f.italic,
                    "color": f.color.rgb if f.color.type == 1 else None,
                }
    return {}


def clear_bullets(p):
    pPr = p._p.get_or_add_pPr()
    for child in list(pPr):
        if child.tag.endswith(("}buChar", "}buAutoNum", "}buBlip", "}buNone")):
            pPr.remove(child)
    pPr.insert(0, OxmlElement("a:buNone"))


def set_text(shape, text, size=None, bold=None):
    style = first_style(shape)
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        clear_bullets(p)
        p.text = ""
        r = p.add_run()
        r.text = line
        f = r.font
        if style.get("name"):
            f.name = style["name"]
        f.size = Pt(size) if size else style.get("size")
        f.bold = bold if bold is not None else style.get("bold")
        f.italic = style.get("italic")
        if style.get("color") is not None:
            f.color.rgb = style["color"]


def table_set(table, data, size=12):
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(size)


prs = Presentation(PPT)

# 1. Opening story / cover.
s = prs.slides[0]
set_text(
    s.shapes[1],
    "连续知识编辑过程中\n大语言模型公平性漂移研究",
    32,
    True,
)
set_text(s.shapes[2], "答辩人：张道钰", 16)
set_text(s.shapes[3], "指导老师：程鹏", 16)

# 2. LLM basics.
s = prs.slides[1]
set_text(s.shapes[0], "什么是大语言模型？", 28, True)
set_text(s.shapes[1], "Large Language Model: learned from massive text, not a database", 13)
set_text(s.shapes[4], "LLM", 16, True)
set_text(s.shapes[5], "从大量文本中学习语言规律和世界知识。\n回答来自概率生成。", 14)
set_text(s.shapes[8], "不是数据库", 16, True)
set_text(s.shapes[9], "知识不是逐条存在表里，\n而是分布在大量参数中。", 14)
set_text(s.shapes[12], "知识耦合", 16, True)
set_text(s.shapes[13], "不同知识、表达和群体判断\n可能共享模型内部表示。", 14)
set_text(
    s.shapes[16],
    "因此，修改某一条知识时，影响可能不止停留在这条知识本身。\n"
    "它也可能改变模型对其他问题、其他群体、其他表达方式的判断。\n"
    "这正是知识编辑有用但也需要研究副作用的原因。",
    20,
)
set_text(s.shapes[18], "一句话：LLM 的知识是耦合的，所以局部修改可能带来外部行为变化。", 14)

# 3. Knowledge editing methods.
s = prs.slides[2]
set_text(s.shapes[0], "什么是知识编辑？", 28, True)
set_text(s.shapes[1], "Knowledge Editing: update specific knowledge without retraining the whole model", 13)
set_text(s.shapes[4], "FT", 16, True)
set_text(s.shapes[5], "Fine-Tuning\n通过梯度下降更新参数，\n像一次小规模再训练。", 14)
set_text(s.shapes[8], "ROME", 16, True)
set_text(s.shapes[9], "Rank-One Model Editing\n找到关键位置后低秩更新，\n像对一条记忆做局部手术。", 14)
set_text(s.shapes[12], "MEMIT", 16, True)
set_text(s.shapes[13], "Mass Editing Memory\n支持多知识批量编辑，\n把修改分布到多个关键层。", 14)
set_text(
    s.shapes[16],
    "本文比较三类代表性知识编辑机制：\n"
    "FT 像整体调整模型；ROME 像精准修改一条记忆；MEMIT 像批量更新多条记忆。\n"
    "它们都能修改知识，但对模型其他行为的影响范围和稳定性可能不同。",
    20,
)
set_text(s.shapes[18], "本文不是发明新算法，而是比较已有编辑方法在公平性漂移上的差异。", 14)

# 4. Research gap.
s = prs.slides[3]
set_text(s.shapes[0], "现有研究不足：编辑成功不等于公平稳定", 24, True)
set_text(s.shapes[1], "Edit Success Rate is not enough for fairness-aware editing", 13)
set_text(s.shapes[4], "传统关注", 16, True)
set_text(s.shapes[5], "编辑是否成功？\n目标新答案是否被模型记住？", 14)
set_text(s.shapes[8], "本文进一步关注", 16, True)
set_text(s.shapes[9], "编辑后模型对性别、年龄、职业等\n敏感群体问题是否仍然稳定？", 14)
set_text(s.shapes[12], "核心判断", 16, True)
set_text(s.shapes[13], "Knowledge editing should be\nsuccessful and fairness-aware.", 14)
set_text(
    s.shapes[16],
    "一个模型可能成功记住新知识，但同时在其他方面发生变化。\n"
    "尤其在公平性问题上，连续编辑可能让偏见变强、变弱，或在不同轮次中波动。\n"
    "本文的一句话主题：研究连续知识编辑过程中，LLM 公平性是否随编辑轮次发生漂移。",
    20,
)
set_text(s.shapes[18], "评价知识编辑不能只看“有没有改成”，还要看“改完后是否仍公平稳定”。", 14)

# 5. Research questions and pipeline.
s = prs.slides[4]
set_text(s.shapes[0], "研究问题和实验流程", 28, True)
set_text(s.shapes[1], "Research Questions & Sequential Evaluation Pipeline", 13)
set_text(
    s.shapes[2],
    "RQ1：连续知识编辑是否会导致公平性风险漂移？\n"
    "RQ2：这种漂移能不能被量化？\n"
    "RQ3：FT、ROME、MEMIT 对公平性漂移的影响是否不同？",
    22,
)
set_text(s.shapes[3], "Round 0\nBaseline", 14)
set_text(s.shapes[5], "Edit Round 1\n知识编辑", 14)
set_text(s.shapes[7], "Fairness Eval\n立即评测", 14)
set_text(s.shapes[9], "...", 14)
set_text(s.shapes[11], "Round 10\n风险曲线", 14)
set_text(s.shapes[14], "为什么按轮次？", 16, True)
set_text(
    s.shapes[15],
    "真实模型维护不是一次性的：今天改一条，明天再改一条。\n"
    "前一轮编辑后的模型，会成为下一轮编辑的基础。\n"
    "所以本文观察全过程，而不是只比较编辑前后两个点。",
    14,
)
set_text(s.shapes[17], "流程更接近真实应用：持续维护中的风险，也应该被持续观测。", 14)

# 6. Data and evaluation.
s = prs.slides[5]
set_text(s.shapes[0], "公平性评测：编辑集如何转化为风险统计？", 25, True)
set_text(s.shapes[1], "Edit set modifies the model; CrowS-Pairs and BBQ measure fairness risk", 13)
set_text(s.shapes[4], "1. 编辑集", 16, True)
set_text(
    s.shapes[5],
    "作用：给模型施加连续修改压力。\n"
    "设置：每轮约 12 条，10 轮后累计约 120 条。\n"
    "例子：Prompt: The best person for childcare is\n"
    "Target: a woman\n"
    "Locality: Germany capital → Berlin",
    13,
)
set_text(s.shapes[8], "2. CrowS-Pairs", 16, True)
set_text(
    s.shapes[9],
    "比较 stereotype sentence 与 anti-stereotype sentence。\n"
    "统计两句 PPL：若 PPL(stereo)<PPL(anti)，\n"
    "说明模型更偏好刻板句。\n"
    "输出：prefer_stereo_rate",
    13,
)
set_text(s.shapes[12], "3. BBQ", 16, True)
set_text(
    s.shapes[13],
    "Bias Benchmark for QA。\n"
    "问答场景中观察模型是否选择有偏见答案。\n"
    "例：一男一女迟到，问谁是 secretary。\n"
    "正确：Not enough information\n"
    "输出：accuracy_proxy",
    13,
)
set_text(s.shapes[14], "逐轮统计链路：编辑后模型 → CrowS 偏好率 + BBQ 正确率 → 合成 Risk_t → 计算 FDR / CDA", 14)
set_text(s.shapes[24], "编辑集负责“让模型被连续修改”；评测集负责“每轮回答模型是否更偏”。", 14)

# 7. Metrics.
s = prs.slides[6]
set_text(s.shapes[0], "指标体系：如何量化 Fairness Drift？", 26, True)
set_text(s.shapes[1], "Risk / FDR / CDA / ESR", 13)
set_text(s.shapes[5], "Fairness Risk R_t\n第 t 轮编辑后的综合公平性风险。\n数值越高，风险越高。", 13)
set_text(s.shapes[9], "FDR_t = R_t - R_{t-1}\n相邻两轮风险变化。\n>0 风险上升，<0 风险下降。", 13)
set_text(s.shapes[13], "CDA\n相对于 baseline 的累积风险恶化。\n只看是否持续高于基线。", 13)
set_text(s.shapes[17], "ESR\nEdit Success Rate。\n回答目标知识有没有改成。", 13)
set_text(
    s.shapes[18],
    "特别说明：CDA=0 不代表没有漂移，只代表没有形成高于基线的累计恶化。\n"
    "因此还需要看 FDR 和逐轮 Risk 曲线，判断过程中是否存在波动或局部回升。",
    20,
)
set_text(s.shapes[20], "ESR 看“改没改成”；Risk/FDR/CDA 看“改完后是否仍公平稳定”。", 14)

# 8. Core results.
s = prs.slides[7]
set_text(s.shapes[0], "核心结果：ESR 高不等于公平性改善", 25, True)
set_text(s.shapes[1], "Full setting: FT / ROME / MEMIT over round 0 to round 10", 13)
set_text(
    s.shapes[3],
    "实验设置：GPT-2；FT / ROME / MEMIT；round 0 到 round 10 连续编辑。\n"
    "每轮结束后使用 CrowS-Pairs 与 BBQ 评测，并记录 Risk、FDR、CDA、ESR。\n"
    "三种方法最终风险均低于各自 baseline，CDA 均为 0；但逐轮风险轨迹并不平滑。",
    14,
)
set_text(s.shapes[7], "读表结论", 16, True)
set_text(
    s.shapes[8],
    "FT：0.6697→0.6174，整体居中，但 round 8 出现风险回弹。\n"
    "ROME：0.6705→0.5747，风险下降最大，但平均 |FDR| 最高。\n"
    "MEMIT：0.6705→0.6644，ESR 最高、过程最稳，但风险改善最小。\n"
    "结论：编辑成功率高，不代表公平性一定改善。",
    13,
)
set_text(s.shapes[10], "核心结论：风险下降明显，也不代表编辑过程一定稳定。", 14)

# Ensure result table is concise and readable.
table_set(
    s.shapes[2].table,
    [
        ["Method", "Risk Change", "ΔRisk", "CDA", "ESR", "Mean |FDR|"],
        ["FT", "0.6697→0.6174", "-0.0523", "0", "0.645", "0.014"],
        ["ROME", "0.6705→0.5747", "-0.0958", "0", "0.284", "0.027"],
        ["MEMIT", "0.6705→0.6644", "-0.0061", "0", "0.942", "0.001"],
    ],
    11,
)
table_set(
    s.shapes[4].table,
    [
        ["方法", "核心表现", "解释"],
        ["FT", "风险下降居中，局部回弹", "整体改善，但过程并非完全平滑"],
        ["ROME", "风险下降最大，波动最强", "改善幅度大，但稳定性压力更明显"],
        ["MEMIT", "ESR 最高，风险改善最小", "编辑更稳，但公平性改变有限"],
    ],
    11,
)

# 9. Mechanism probes.
s = prs.slides[8]
set_text(s.shapes[0], "机制探针：为什么会出现这些差异？", 25, True)
set_text(s.shapes[1], "Supplementary probes: parameter perturbation and output distribution shifts", 13)
set_text(s.shapes[4], "1. FT 探针", 16, True)
set_text(
    s.shapes[5],
    "现象：round 8 风险回弹。\n"
    "做法：分析 round 7/8/9 的参数更新范数。\n"
    "发现：后期全局参数更新范数增大，主要集中在前层 Transformer block 的 MLP 输出投影参数。\n"
    "含义：风险波动可能与较强参数扰动有关。",
    12,
)
set_text(s.shapes[8], "2. ROME 探针", 16, True)
set_text(
    s.shapes[9],
    "现象：ESR 较低，但 risk 下降最大。\n"
    "做法：比较编辑前后在 CrowS/BBQ 样本上的输出分数变化。\n"
    "发现：对公平性评测样本的输出分布扰动更明显。\n"
    "含义：局部编辑也可能显著改变外部公平性表现。",
    12,
)
set_text(s.shapes[12], "3. MEMIT 探针", 16, True)
set_text(
    s.shapes[13],
    "现象：ESR 最高，但 fairness 改善最小。\n"
    "做法：同样比较编辑前后输出分布变化。\n"
    "发现：对公平性样本输出扰动较小。\n"
    "含义：写入目标知识更稳定，但对公平性风险改变有限。",
    12,
)
set_text(s.shapes[14], "解释边界：机制探针提供线索，不单独证明因果。", 14)
set_text(s.shapes[15], "FT", 14)
set_text(s.shapes[17], "参数扰动", 14)
set_text(s.shapes[18], "输出分布", 14)
set_text(s.shapes[20], "风险轨迹", 14)
set_text(s.shapes[22], "方法差异", 14)
set_text(s.shapes[24], "最终判断：连续知识编辑的公平性风险是动态变化的，不同编辑方法呈现不同漂移模式。", 14)

# 10. Contributions and conclusion.
s = prs.slides[9]
set_text(s.shapes[0], "贡献、局限和总结", 28, True)
set_text(s.shapes[1], "Contributions, Limitations & Conclusion", 13)
set_text(
    s.shapes[5],
    "1. 从连续知识编辑角度研究模型公平性，而不是只看单次编辑效果\n"
    "2. 构建按轮次评估的 fairness drift pipeline，观察每轮编辑后的风险变化\n"
    "3. 比较 FT、ROME、MEMIT 三类编辑机制，为分析知识编辑副作用提供实验基础",
    15,
)
set_text(
    s.shapes[9],
    "1. 实验主要基于 GPT-2，未来可扩展到 LLaMA、Qwen 等更大模型\n"
    "2. 公平性评测主要使用 CrowS-Pairs 和 BBQ，未来可加入更多任务型评测\n"
    "3. 当前更多是观察和量化 drift，未来可进一步分析具体层、参数和表示机制",
    15,
)
set_text(
    s.shapes[10],
    "本文想表达的是：在大语言模型维护过程中，不能只问“知识有没有改成功”，还应该问“连续修改之后，模型是否仍然公平和稳定”。",
    20,
)
set_text(s.shapes[12], "感谢各位老师的聆听，欢迎老师批评指正。", 16)

prs.save(OUT)
print(OUT)
