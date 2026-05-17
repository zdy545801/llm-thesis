from pathlib import Path

from pptx import Presentation
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt


PPT = Path(r"D:\学习\222_讲稿优化11页最终版.pptx")


def first_style(shape):
    if not hasattr(shape, "text_frame"):
        return {}
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if r.text:
                f = r.font
                return {
                    "name": f.name,
                    "size": f.size,
                    "bold": f.bold,
                    "italic": f.italic,
                    "color": f.color.rgb if f.color.type == 1 else None,
                }
    return {}


def set_text(shape, text, size=None, bold=None):
    style = first_style(shape)
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        pPr = p._p.get_or_add_pPr()
        for child in list(pPr):
            if child.tag.endswith(("}buChar", "}buAutoNum", "}buBlip", "}buNone")):
                pPr.remove(child)
        pPr.insert(0, OxmlElement("a:buNone"))
        p.text = ""
        r = p.add_run()
        r.text = line
        f = r.font
        if style.get("name"):
            f.name = style["name"]
        f.size = Pt(size) if size else style.get("size")
        f.bold = bold if bold is not None else style.get("bold")
        f.italic = style.get("italic")
        if style.get("color") is not None:
            f.color.rgb = style["color"]


prs = Presentation(PPT)
s = prs.slides[1]
set_text(s.shapes[0], "开头故事：模型被不断修正后还公平吗？", 26, True)
set_text(s.shapes[1], "Motivation: continuous maintenance may create hidden fairness side effects", 13)
set_text(s.shapes[4], "真实系统上线", 16, True)
set_text(s.shapes[5], "在线客服、教育问答、招聘辅助等系统中，\n模型会持续面对新需求。", 14)
set_text(s.shapes[8], "发现知识错误", 16, True)
set_text(s.shapes[9], "人物职位变化、政策更新、事实过时，\n不可能每次都重训整个大模型。", 14)
set_text(s.shapes[12], "连续知识编辑", 16, True)
set_text(s.shapes[13], "今天改一条，明天再改一条。\n多轮修改后可能出现未预料副作用。", 14)
set_text(
    s.shapes[16],
    "知识编辑（Knowledge Editing）让我们能在不重新训练整个模型的情况下，直接修改模型内部的某些知识。\n"
    "但如果模型原本已经存在偏见，连续编辑后这些偏见可能变强、变弱，或者在不同轮次中波动。\n"
    "这就是本文关注的公平性漂移（Fairness Drift）。",
    20,
)
set_text(s.shapes[18], "研究动机：持续维护模型时，不能只问“有没有改对”，还要问“公平性是否仍稳定”。", 14)

prs.save(PPT)
print(PPT)
