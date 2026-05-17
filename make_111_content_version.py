from pathlib import Path
import zipfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


SRC = Path(r"D:\学习\111.pptx")
OUT = Path(r"D:\学习\111_答辩内容版.pptx")
LOGO = Path(r"D:\毕业论文LLM\tmp_ppt_media\111_image1.png")

BLUE = RGBColor(37, 51, 129)
BLUE2 = RGBColor(56, 93, 198)
CYAN = RGBColor(58, 160, 194)
LIGHT_BLUE = RGBColor(236, 246, 255)
MID_BLUE = RGBColor(190, 214, 246)
TEXT = RGBColor(34, 39, 52)
GRAY = RGBColor(113, 120, 132)
LIGHT_GRAY = RGBColor(246, 248, 252)
WHITE = RGBColor(255, 255, 255)
RED = RGBColor(190, 73, 64)


def ensure_logo():
    if LOGO.exists():
        return
    LOGO.parent.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(SRC) as z:
        LOGO.write_bytes(z.read("ppt/media/image1.png"))


def set_run(run, size=16, bold=False, color=TEXT, font="Microsoft YaHei"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(slide, x, y, w, h, text, size=16, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_run(r, size=size, bold=bold, color=color)
    return box


def add_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = text


def add_header(slide, title):
    # Left title marker: two triangles, matching the provided content-page style.
    tri1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_TRIANGLE, Inches(0.65), Inches(0.55), Inches(0.34), Inches(0.34))
    tri1.fill.solid()
    tri1.fill.fore_color.rgb = BLUE2
    tri1.line.fill.background()
    tri2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_TRIANGLE, Inches(0.84), Inches(0.55), Inches(0.34), Inches(0.34))
    tri2.fill.solid()
    tri2.fill.fore_color.rgb = CYAN
    tri2.line.fill.background()
    add_text(slide, 1.23, 0.45, 6.6, 0.55, title, size=27, bold=True, color=BLUE)
    slide.shapes.add_picture(str(LOGO), Inches(10.35), Inches(0.48), width=Inches(2.25))
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.65), Inches(1.13), Inches(12.0), Inches(0.012))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(205, 205, 205)
    line.line.fill.background()
    foot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.65), Inches(6.46), Inches(12.0), Inches(0.012))
    foot.fill.solid()
    foot.fill.fore_color.rgb = RGBColor(205, 205, 205)
    foot.line.fill.background()
    add_text(slide, 0.68, 6.6, 4.2, 0.28, "含弘光大   继往开来", size=14.5, color=BLUE2)
    add_text(slide, 10.35, 6.6, 2.3, 0.28, "www.swu.edu.cn", size=14.5, color=BLUE2, align=PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, title, body, accent=BLUE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = MID_BLUE
    shape.line.width = Pt(1.0)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_text(slide, x + 0.2, y + 0.14, w - 0.35, 0.28, title, size=15.5, bold=True, color=BLUE)
    add_text(slide, x + 0.2, y + 0.52, w - 0.35, h - 0.6, body, size=12.4, color=TEXT)
    return shape


def add_bullets(slide, x, y, w, h, bullets, size=17):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(8)
    return box


def add_table(slide, x, y, w, h, data, widths=None, font_size=10.5, header=BLUE):
    ts = slide.shapes.add_table(len(data), len(data[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = ts.table
    if widths:
        for i, cw in enumerate(widths):
            table.columns[i].width = Inches(cw)
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.margin_left = Inches(0.03)
            cell.margin_right = Inches(0.03)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header if r == 0 else (LIGHT_GRAY if r % 2 == 0 else WHITE)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c != 0 else PP_ALIGN.LEFT
            for run in p.runs:
                set_run(run, size=font_size, bold=(r == 0), color=WHITE if r == 0 else TEXT)
    return ts


def add_label(slide, x, y, w, h, text, fill=LIGHT_BLUE, color=BLUE, size=13):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = MID_BLUE
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    set_run(r, size=size, bold=True, color=color)


def add_arrow(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = BLUE2
    line.line.width = Pt(1.8)
    line.line.end_arrowhead = True


def remove_slides_after_first(prs):
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst)[1:]:
        prs.part.drop_rel(sld_id.rId)
        sld_id_lst.remove(sld_id)


def make_deck():
    ensure_logo()
    prs = Presentation(str(SRC))
    remove_slides_after_first(prs)
    blank = prs.slide_layouts[6]
    add_notes(prs.slides[0], "各位老师好，我是张道钰。我的论文题目是《大语言模型知识编辑中的公平性漂移研究》。接下来我将从研究背景、研究框架、数据指标、实验结果和结论展望几个方面进行汇报。")

    # 2. Directory
    s = prs.slides.add_slide(blank)
    add_header(s, "目录")
    add_text(s, 1.15, 1.7, 10.5, 0.35, "答辩内容结构", size=22, bold=True, color=BLUE)
    items = [
        ("01", "研究背景与问题提出", "为什么关注连续编辑中的公平性风险"),
        ("02", "研究框架与数据评测", "怎样逐轮编辑、逐轮测量"),
        ("03", "指标体系与实验结果", "Risk / FDR / CDA / ESR 得到什么结论"),
        ("04", "贡献、局限与总结", "本文价值和后续方向"),
    ]
    for i, (num, title, desc) in enumerate(items):
        y = 2.25 + i * 0.82
        add_label(s, 1.3, y, 0.78, 0.42, num, fill=BLUE, color=WHITE, size=12.5)
        add_text(s, 2.28, y - 0.02, 3.8, 0.34, title, size=17.5, bold=True, color=TEXT)
        add_text(s, 6.1, y + 0.03, 5.4, 0.28, desc, size=13, color=GRAY)
    add_notes(s, "本次答辩分为四部分：研究背景与问题提出，研究框架与数据评测，指标体系与实验结果，最后是贡献、局限和总结。")

    # 3. Background
    s = prs.slides.add_slide(blank)
    add_header(s, "研究背景：模型被不断修正后还公平吗？")
    add_card(s, 0.9, 1.55, 3.45, 1.25, "LLM", "知识分布在大量参数中，不是逐条存在数据库里。", accent=BLUE)
    add_card(s, 4.95, 1.55, 3.45, 1.25, "Knowledge Editing", "不重训整个模型，定向修改局部知识。", accent=BLUE2)
    add_card(s, 9.0, 1.55, 3.45, 1.25, "Sequential Editing", "在同一模型上连续进行多轮修改。", accent=RED)
    add_arrow(s, 4.42, 2.18, 4.87, 2.18)
    add_arrow(s, 8.48, 2.18, 8.9, 2.18)
    add_bullets(s, 1.15, 3.35, 10.9, 1.45, [
        "真实系统上线后，模型会随着新需求、错误反馈和风险发现不断被修正",
        "单次编辑只回答“这一次有没有改成”，连续编辑还要回答“多轮之后是否仍稳定”",
        "本文关注：连续编辑是否会引发公平性漂移（Fairness Drift）"
    ], size=17)
    add_notes(s, "这一页解释研究动机。大语言模型上线后会被不断修正。知识编辑可以在不重新训练整个模型的情况下修改局部知识。但连续编辑和单次编辑不同，前一轮编辑后的模型会成为后一轮编辑的基础，因此可能带来累积影响。本文关注连续修改后模型公平性是否仍然稳定。")

    # 4. Framework
    s = prs.slides.add_slide(blank)
    add_header(s, "研究问题与总体框架")
    add_bullets(s, 0.95, 1.52, 5.25, 2.2, [
        "连续知识编辑是否会引发公平性漂移？",
        "如何量化风险的轮间变化和累计漂移？",
        "FT、ROME、MEMIT 的风险轨迹是否不同？",
        "编辑成功率 ESR 是否意味着公平性同步改善？"
    ], size=16)
    steps = ["Round 0 基线评测", "Edit Round 1", "Fairness Eval", "...", "Round 10 分析"]
    for i, lab in enumerate(steps):
        y = 1.52 + i * 0.72
        add_label(s, 7.25, y, 3.6, 0.42, lab, fill=WHITE if i % 2 == 0 else LIGHT_BLUE, color=BLUE, size=12.5)
        if i < len(steps) - 1:
            add_arrow(s, 9.05, y + 0.43, 9.05, y + 0.68)
    add_card(s, 0.95, 4.62, 5.25, 0.72, "实验对象", "GPT-2 + FT / ROME / MEMIT；连续 10 轮逐轮评测", accent=BLUE)
    add_notes(s, "本文围绕四个问题展开：是否发生公平性漂移，如何量化这种漂移，不同方法之间是否存在差异，以及 ESR 是否意味着公平性同步改善。实验流程是先得到 round 0 的基线风险，然后每轮编辑后都进行公平性评测，一直到 round 10。")

    # 5. Data page
    s = prs.slides.add_slide(blank)
    add_header(s, "数据与评测：编辑集如何转化为风险统计？")
    add_card(s, 0.85, 1.35, 3.75, 2.7, "1. 编辑集：施加连续修改", "作用：每轮累积编辑请求。\n\n例子：\nPrompt: The best person for childcare is\nTarget: a woman\nLocality: Germany capital → Berlin\n\n本实验：10 轮，最终 120 条。", accent=BLUE)
    add_card(s, 4.8, 1.35, 3.75, 2.7, "2. CrowS-Pairs：测刻板偏好", "形式：刻板句 vs 反刻板句。\n\n统计：比较两句 PPL。\n若 PPL(stereo)<PPL(anti)，记为偏好刻板句。\n\n输出：prefer_stereo_rate", accent=BLUE2)
    add_card(s, 8.75, 1.35, 3.75, 2.7, "3. BBQ：测问答偏见", "例子：一男一女迟到，问谁是 secretary。\n正确：Not enough information。\n\n统计：选择 PPL 最低答案。\n输出：accuracy_proxy", accent=RED)
    add_text(s, 1.05, 4.55, 11.1, 0.24, "逐轮统计：编辑后模型 → CrowS/BBQ 评测 → 综合风险 Risk_t → FDR / CDA", size=15.5, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_label(s, 1.45, 5.08, 1.85, 0.4, "编辑后模型", fill=WHITE, color=BLUE)
    add_arrow(s, 3.38, 5.28, 3.85, 5.28)
    add_label(s, 4.0, 4.86, 2.18, 0.38, "CrowS 偏好率", fill=LIGHT_BLUE, color=BLUE, size=11)
    add_label(s, 4.0, 5.56, 2.18, 0.38, "BBQ 正确率", fill=LIGHT_BLUE, color=BLUE, size=11)
    add_arrow(s, 6.25, 5.28, 6.72, 5.28)
    add_label(s, 6.85, 5.08, 3.2, 0.4, "Risk_t = 0.5×CrowS + 0.5×(1-BBQ)", fill=BLUE, color=WHITE, size=10.5)
    add_arrow(s, 10.15, 5.28, 10.6, 5.28)
    add_label(s, 10.75, 5.08, 1.2, 0.4, "FDR/CDA", fill=WHITE, color=BLUE, size=11)
    add_notes(s, "这一页解释数据怎么进入实验。编辑集给模型施加连续修改压力。每轮编辑后，用 CrowS-Pairs 和 BBQ 评测。CrowS 比较刻板句和反刻板句的困惑度，统计模型更偏好刻板句的比例。BBQ 是问答评测，正确答案通常是信息不足，如果模型选择某个群体就可能体现偏见。最后把 CrowS 偏好率和 BBQ 错误风险合成 risk，再逐轮计算 FDR 和 CDA。")

    # 6. Metrics and methods
    s = prs.slides.add_slide(blank)
    add_header(s, "指标体系：如何量化 Fairness Drift？")
    add_card(s, 0.82, 1.42, 2.9, 1.8, "Risk_t", "第 t 轮综合公平性风险。\n数值越高，风险越高。", accent=BLUE)
    add_card(s, 3.85, 1.42, 2.9, 1.8, "FDR_t", "FDR_t = Risk_t - Risk_t-1\n表示相邻两轮风险变化。", accent=BLUE2)
    add_card(s, 6.88, 1.42, 2.9, 1.8, "CDA", "累计高于初始基线的风险恶化。\nCDA=0 不代表没有漂移。", accent=CYAN)
    add_card(s, 9.9, 1.42, 2.25, 1.8, "ESR", "目标知识是否编辑成功。\n高 ESR 不一定代表风险降低。", accent=RED)
    add_bullets(s, 1.05, 4.02, 11.05, 1.12, [
        "ESR 回答“是否改成”；Risk / FDR / CDA 回答“改后公平性如何变化”",
        "FDR 捕捉轮间波动，CDA 捕捉相对初始基线的累计不利漂移",
        "多指标联合分析避免只依据单一数值判断方法优劣"
    ], size=16)
    add_notes(s, "这一页解释指标。Risk 表示每一轮后的综合公平性风险。FDR 表示相邻两轮之间风险变化，大于 0 代表风险上升，小于 0 代表风险下降。CDA 只累计高于初始基线的风险恶化面积，所以 CDA 为 0 不等于模型没有变化，而是没有形成高于基线的累计恶化。ESR 衡量目标知识是否编辑成功。")

    # 7. Results
    s = prs.slides.add_slide(blank)
    add_header(s, "实验结果、核心结论与机制探针")
    result = [
        ["Method", "Risk Change", "ΔRisk", "CDA", "ESR", "Mean |FDR|"],
        ["FT", "0.6697→0.6174", "-0.0523", "0", "0.645", "0.014"],
        ["ROME", "0.6705→0.5747", "-0.0958", "0", "0.284", "0.027"],
        ["MEMIT", "0.6705→0.6644", "-0.0061", "0", "0.942", "0.001"],
    ]
    add_table(s, 0.78, 1.35, 7.25, 1.45, result, widths=[1.0, 1.75, 0.9, 0.58, 0.8, 1.25], font_size=9.8)
    add_bullets(s, 8.35, 1.25, 3.8, 1.55, [
        "ROME：风险下降最大，但波动最强",
        "MEMIT：ESR 最高、波动最小，但改善最弱",
        "ESR 高 ≠ 公平性一定改善"
    ], size=13)
    probe = [
        ["方法", "探针结果", "解释"],
        ["FT", "后期参数更新范数增大", "风险回弹可能与较强扰动有关"],
        ["ROME", "输出分布扰动最大", "低 ESR 但 risk 降幅大"],
        ["MEMIT", "输出分布扰动最小", "高 ESR、低波动、改善有限"],
    ]
    add_table(s, 0.78, 3.68, 7.25, 1.42, probe, widths=[0.9, 2.25, 4.1], font_size=9.8, header=RGBColor(56, 91, 133))
    add_card(s, 8.35, 3.68, 3.8, 1.42, "最终判断", "连续编辑中的公平性风险是动态变化的；不同方法在风险下降、波动强度和编辑成功率之间存在取向差异。", accent=RED)
    add_notes(s, "这一页是核心结果。FT 风险从 0.6697 降到 0.6174，ROME 从 0.6705 降到 0.5747，MEMIT 从 0.6705 降到 0.6644。三种方法 CDA 都是 0，说明没有形成高于初始基线的累计恶化，但逐轮 FDR 显示风险仍在波动。ROME 风险下降最大，但波动也最强；MEMIT ESR 最高，但公平性改善最小。机制探针进一步说明，FT 后期参数更新范数增大，ROME 对公平性样本输出分布扰动明显，MEMIT 扰动较小。")

    # 8. Conclusion
    s = prs.slides.add_slide(blank)
    add_header(s, "贡献、局限与总结")
    add_card(s, 0.95, 1.45, 5.45, 2.0, "主要贡献", "1. 从连续编辑角度研究公平性风险\n2. 构建逐轮评测框架，量化 Risk / FDR / CDA\n3. 比较 FT、ROME、MEMIT 的风险轨迹\n4. 发现 ESR 与 fairness improvement 不一定一致", accent=BLUE)
    add_card(s, 6.95, 1.45, 5.45, 2.0, "局限与展望", "1. 主要基于 GPT-2，尚未扩展到更大模型\n2. 评测主要依赖 CrowS-Pairs 和 BBQ\n3. 机制探针仍是小规模解释性证据\n4. 后续需要更多随机种子、更长编辑链路", accent=BLUE2)
    add_label(s, 1.4, 4.65, 10.6, 0.62, "知识编辑不能只关注“是否改对”，还应关注连续修改中模型是否仍然稳定、公平和可控。", fill=LIGHT_BLUE, color=BLUE, size=16)
    add_notes(s, "最后总结本文贡献和局限。本文从连续编辑角度研究公平性风险，构建逐轮评测框架，比较三种代表性方法，并发现 ESR 和公平性改善不一定同步。局限是主要基于 GPT-2，评测集主要是 CrowS-Pairs 和 BBQ，机制探针还是小规模解释性证据。未来可以扩展到更大模型、更多评测维度、更多随机种子和更长编辑链路。")

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    p = make_deck()
    print(f"saved: {p}")
