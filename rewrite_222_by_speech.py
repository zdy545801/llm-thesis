from pathlib import Path

from pptx import Presentation
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt


SRC = Path(r"D:\学习\222_内容丰富版.pptx")
DST = Path(r"D:\学习\222_按讲稿重写版.pptx")


def _first_font(shape):
    if not hasattr(shape, "text_frame"):
        return {}
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if r.text:
                f = r.font
                return {
                    "name": f.name,
                    "size": f.size,
                    "bold": f.bold,
                    "italic": f.italic,
                    "color": f.color.rgb if f.color.type == 1 else None,
                }
    return {}


def set_text(shape, text, size=None, bold=None):
    style = _first_font(shape)
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        pPr = p._p.get_or_add_pPr()
        for child in list(pPr):
            if child.tag.endswith(("}buChar", "}buAutoNum", "}buBlip", "}buNone")):
                pPr.remove(child)
        pPr.insert(0, OxmlElement("a:buNone"))
        p.text = ""
        r = p.add_run()
        r.text = line
        f = r.font
        if style.get("name"):
            f.name = style["name"]
        f.size = Pt(size) if size else style.get("size")
        f.bold = bold if bold is not None else style.get("bold")
        f.italic = style.get("italic")
        if style.get("color") is not None:
            f.color.rgb = style["color"]


prs = Presentation(SRC)

# 1. Opening story / cover.
s = prs.slides[0]
set_text(
    s.shapes[1],
    "连续知识编辑过程中\n大语言模型公平性漂移研究\nFairness Drift in Continuous Knowledge Editing of LLMs",
    31,
    True,
)
set_text(s.shapes[2], "答辩人：张道钰", 16)
set_text(s.shapes[3], "指导老师：程鹏", 16)

# 2. What is LLM? Reuse the directory page without moving shapes.
s = prs.slides[1]
set_text(s.shapes[1], "什么是大语言模型（LLM）？", 22, True)
set_text(s.shapes[0], "A model learned from massive text", 13)
set_text(s.shapes[3], "知识不是数据库式存放", 22, True)
set_text(s.shapes[2], "Knowledge is distributed in parameters", 13)
set_text(s.shapes[5], "知识之间存在耦合", 22, True)
set_text(s.shapes[4], "Changing one part may affect others", 13)
set_text(s.shapes[7], "这就是副作用风险", 22, True)
set_text(s.shapes[6], "Editing is useful, but its side effects matter", 13)

# 3. What is knowledge editing?
s = prs.slides[2]
set_text(s.shapes[0], "什么是知识编辑？", 26, True)
set_text(s.shapes[1], "Knowledge Editing: update specific knowledge without retraining the whole model", 13)
set_text(s.shapes[4], "FT", 16, True)
set_text(s.shapes[5], "Fine-Tuning：通过梯度下降更新参数。\n像“小规模再训练”，直接但影响范围可能较大。", 14)
set_text(s.shapes[8], "ROME", 16, True)
set_text(s.shapes[9], "Rank-One Model Editing：定位关键位置后低秩更新。\n像“局部手术”，用于单条知识精确修改。", 14)
set_text(s.shapes[12], "MEMIT", 16, True)
set_text(s.shapes[13], "Mass Editing Memory in a Transformer：支持多知识批量编辑。\n像“批量更新多条记忆”，修改分布到多个关键层。", 14)
set_text(
    s.shapes[16],
    "本文比较三类代表性编辑机制：FT 代表整体参数调整；ROME 代表单知识局部编辑；MEMIT 代表多知识批量编辑。\n"
    "这三类方法都能修改模型知识，但它们对模型其他行为的影响范围和稳定性可能不同。",
    18,
)
set_text(s.shapes[18], "简单理解：FT 像整体调整模型，ROME 像精准修改一条记忆，MEMIT 像批量更新多条记忆。", 14)

# 4. Research gap.
s = prs.slides[3]
set_text(s.shapes[0], "现有研究不足：编辑成功不等于公平稳定", 24, True)
set_text(s.shapes[1], "Edit Success ≠ Fairness Stability", 14)
set_text(
    s.shapes[2],
    "现有知识编辑研究常问：编辑是否成功？\n"
    "也就是模型是否记住了目标新答案，通常用 ESR 衡量。\n"
    "但一个模型可能“改对了目标知识”，同时在性别、年龄、种族、职业等敏感群体问题上产生新的偏向。\n"
    "因此，Knowledge editing should not only be successful, but also fairness-aware.",
    20,
)
set_text(s.shapes[3], "传统关注\nEdit Success", 14)
set_text(s.shapes[5], "本文关注\nFairness-Aware", 14)
set_text(s.shapes[7], "核心问题\n多轮后是否稳定", 14)
set_text(s.shapes[9], "动态观察\n逐轮风险曲线", 14)
set_text(s.shapes[11], "本文主题\nFairness Drift", 14)
set_text(
    s.shapes[15],
    "一句话主题：研究“连续知识编辑”过程中，LLM 公平性是否随编辑轮次发生漂移，并给出可复现的量化指标与实验流程。",
    14,
)
set_text(s.shapes[17], "知识编辑不能只回答“有没有改成功”，还要回答“改完后模型是否仍公平稳定”。", 14)

# 5. Research questions and experiment flow.
s = prs.slides[4]
set_text(s.shapes[0], "研究问题和实验流程", 26, True)
set_text(s.shapes[1], "Research Questions & Sequential Evaluation Pipeline", 13)
set_text(
    s.shapes[5],
    "RQ1：连续知识编辑是否导致公平性风险漂移？\n"
    "RQ2：这种漂移能不能被量化？\n"
    "RQ3：FT、ROME、MEMIT 对公平性漂移的影响是否不同？\n"
    "目标：不只比较编辑前后两个点，而是观察风险曲线如何变化。",
    15,
)
set_text(
    s.shapes[9],
    "实验流程：\n"
    "1. 初始模型先做公平性评测，得到 baseline\n"
    "2. 进行第 1 轮知识编辑\n"
    "3. 编辑后马上再次评测公平性\n"
    "4. 继续第 2 轮、第 3 轮……直到 round 10",
    14,
)
set_text(
    s.shapes[13],
    "为什么按轮次？\n真实模型维护往往不是一次性的。\n今天改一条、明天再改一条，前一轮编辑后的模型会成为下一轮基础。",
    14,
)
set_text(s.shapes[14], "逐轮统计链路：baseline → edit → fairness eval → edit → fairness eval → risk curve", 14)
set_text(s.shapes[15], "Round 0\nBaseline", 14)
set_text(s.shapes[17], "Edit + Eval\nRound 1...10", 14)
set_text(s.shapes[18], "Risk Curve\nFDR / CDA", 14)
set_text(s.shapes[20], "Compare\nFT / ROME / MEMIT", 14)
set_text(s.shapes[22], "Conclusion\nDrift Pattern", 14)
set_text(s.shapes[24], "流程更接近真实场景：模型维护是持续发生的，风险也应逐轮观察。", 14)

# 6. Evaluation datasets and metrics.
s = prs.slides[5]
set_text(s.shapes[0], "公平性评测和指标", 26, True)
set_text(s.shapes[1], "CrowS-Pairs / BBQ + Risk / FDR / CDA", 13)
set_text(
    s.shapes[5],
    "CrowS-Pairs：\n比较 stereotype sentence 与 anti-stereotype sentence。\n如果模型更偏好刻板句，说明刻板偏好风险更高。",
    13,
)
set_text(
    s.shapes[9],
    "BBQ：Bias Benchmark for QA。\n通过问答场景观察模型是否在敏感群体问题中选择有偏见答案。",
    13,
)
set_text(
    s.shapes[13],
    "Fairness Risk R_t：\n第 t 轮编辑后模型在公平性评测集上的综合风险。\n数值越高，风险越高。",
    13,
)
set_text(
    s.shapes[17],
    "FDR_t = R_t - R_{t-1}\n大于 0：本轮后风险上升。\n小于 0：本轮后风险下降。\nCDA：相对基线的累积风险恶化。",
    12,
)
set_text(
    s.shapes[18],
    "指标解释：FDR 观察相邻两轮之间的风险变化；CDA 观察多轮编辑后是否形成高于基线的累计恶化。\n"
    "特别注意：CDA=0 不代表没有漂移，只代表没有形成高于基线的累积恶化。",
    20,
)
set_text(s.shapes[20], "本文用 CrowS-Pairs 和 BBQ 评估公平性，用 Risk / FDR / CDA 描述漂移过程。", 14)

# 7. Experimental setting, core results, and mechanism probes.
s = prs.slides[6]
set_text(s.shapes[0], "实验设置、核心结果与机制探针", 26, True)
set_text(s.shapes[1], "Experimental Setting, Main Results & Mechanism Probes", 13)
set_text(
    s.shapes[3],
    "实验设置：GPT-2；FT / ROME / MEMIT；round 0 到 round 10 连续编辑。\n"
    "每轮编辑后使用 CrowS-Pairs 与 BBQ 评测，并记录 Risk、FDR、CDA、ESR。\n\n"
    "Full 结果：\n"
    "FT：0.6697→0.6174，Δ=-0.0523，ESR=0.645\n"
    "ROME：0.6705→0.5747，Δ=-0.0958，ESR=0.284\n"
    "MEMIT：0.6705→0.6644，Δ=-0.0061，ESR=0.942\n"
    "三者 CDA 均为 0，但逐轮风险轨迹并不平滑。",
    13,
)
set_text(
    s.shapes[8],
    "核心结论：ESR 高不代表公平性一定改善；风险下降明显，也不代表编辑过程一定稳定。\n"
    "FT：整体居中，但 round 8 出现风险回弹。\n"
    "ROME：风险压降最大，但平均 |FDR| 最高，波动最强。\n"
    "MEMIT：ESR 最高、过程最平稳，但风险改善最小。\n\n"
    "机制探针：FT 后期参数更新范数增大；ROME 输出分布扰动更明显；MEMIT 输出扰动更小。",
    13,
)
set_text(
    s.shapes[10],
    "最终判断：连续知识编辑中的公平性风险是动态变化的；不同方法呈现不同漂移模式。",
    13,
)

# 8. Contributions, limitations, and conclusion.
s = prs.slides[7]
set_text(s.shapes[0], "贡献、局限和总结", 26, True)
set_text(s.shapes[1], "Contributions, Limitations & Conclusion", 13)
set_text(
    s.shapes[5],
    "1. 从连续知识编辑角度研究模型公平性，而不是只看单次编辑效果\n"
    "2. 构建按轮次评估的 fairness drift pipeline，观察每轮编辑后的风险变化\n"
    "3. 比较 FT、ROME、MEMIT 三类不同编辑机制，为分析知识编辑副作用提供实验基础",
    15,
)
set_text(
    s.shapes[9],
    "1. 实验主要基于 GPT-2，未来可扩展到 LLaMA、Qwen 等更大模型\n"
    "2. 公平性评测主要使用 CrowS-Pairs 和 BBQ，未来可加入更多任务型评测\n"
    "3. 当前更多是观察和量化 drift，未来可进一步分析具体层、参数和表示机制",
    15,
)
set_text(
    s.shapes[10],
    "本文想表达的是：在大语言模型维护过程中，不能只问“知识有没有改成功”，还应该问“连续修改之后，模型是否仍然公平和稳定”。",
    20,
)
set_text(s.shapes[12], "感谢各位老师的聆听，欢迎老师批评指正。", 16)

prs.save(DST)
print(DST)
