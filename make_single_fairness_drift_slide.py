from pathlib import Path

from pptx import Presentation
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt


PPT = Path(r"D:\学习\单页_现实场景到公平性漂移.pptx")


def first_style(shape):
    if not hasattr(shape, "text_frame"):
        return {}
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if r.text:
                f = r.font
                return {
                    "name": f.name,
                    "size": f.size,
                    "bold": f.bold,
                    "italic": f.italic,
                    "color": f.color.rgb if f.color.type == 1 else None,
                }
    return {}


def clear_bullets(p):
    pPr = p._p.get_or_add_pPr()
    for child in list(pPr):
        if child.tag.endswith(("}buChar", "}buAutoNum", "}buBlip", "}buNone")):
            pPr.remove(child)
    pPr.insert(0, OxmlElement("a:buNone"))


def set_text(shape, text, size=None, bold=None):
    style = first_style(shape)
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        clear_bullets(p)
        p.text = ""
        run = p.add_run()
        run.text = line
        f = run.font
        if style.get("name"):
            f.name = style["name"]
        f.size = Pt(size) if size else style.get("size")
        f.bold = bold if bold is not None else style.get("bold")
        f.italic = style.get("italic")
        if style.get("color") is not None:
            f.color.rgb = style["color"]


prs = Presentation(PPT)
s = prs.slides[0]

set_text(s.shapes[0], "现实场景：为什么需要研究公平性漂移？", 25, True)
set_text(s.shapes[1], "From practical model maintenance to fairness drift", 13)

set_text(s.shapes[4], "1. 现实场景", 16, True)
set_text(
    s.shapes[5],
    "LLM 已应用于在线客服、教育问答、招聘辅助等系统。\n"
    "上线后可能出现：\n"
    "知识错误 / 知识过时 / 回答不合理或带偏见",
    13,
)

set_text(s.shapes[8], "2. 为什么需要知识编辑？", 16, True)
set_text(
    s.shapes[9],
    "每次发现问题都重新训练整个模型，成本太高。\n"
    "知识编辑（Knowledge Editing）：\n"
    "不重训全模型，对局部知识进行定向修改。",
    13,
)

set_text(s.shapes[12], "3. 核心问题", 16, True)
set_text(
    s.shapes[13],
    "真实应用中，模型往往会被连续、多轮修改。\n"
    "编辑本质上改变下一个词的概率预测；\n"
    "模型并不知道自己的回答是否带有偏见。",
    13,
)

set_text(
    s.shapes[16],
    "本文关注：连续知识编辑后，模型的公平性风险会不会发生变化？\n"
    "这种变化可能表现为增强、减弱或波动。\n"
    "本文将其称为公平性漂移（Fairness Drift）。",
    24,
)

set_text(s.shapes[18], "核心引出：知识编辑解决“改模型”，公平性漂移关注“连续修改后是否仍公平、稳定”。", 14)

prs.save(PPT)
print(PPT)
