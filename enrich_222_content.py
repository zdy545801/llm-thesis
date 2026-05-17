from pathlib import Path

from pptx import Presentation
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt


SRC = Path(r"D:\学习\222.pptx")
DST = Path(r"D:\学习\222_内容丰富版.pptx")


def _first_font(shape):
    """Capture the first visible run's style so replacing text does not restyle the slide."""
    if not hasattr(shape, "text_frame"):
        return {}
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if r.text:
                f = r.font
                return {
                    "name": f.name,
                    "size": f.size,
                    "bold": f.bold,
                    "italic": f.italic,
                    "color": f.color.rgb if f.color.type == 1 else None,
                }
    return {}


def set_text(shape, text, size=None, bold=None):
    style = _first_font(shape)
    tf = shape.text_frame
    tf.clear()
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        pPr = p._p.get_or_add_pPr()
        for child in list(pPr):
            if child.tag.endswith(("}buChar", "}buAutoNum", "}buBlip", "}buNone")):
                pPr.remove(child)
        pPr.insert(0, OxmlElement("a:buNone"))
        p.text = ""
        run = p.add_run()
        run.text = line
        font = run.font
        if style.get("name"):
            font.name = style["name"]
        font.size = Pt(size) if size else style.get("size")
        font.bold = bold if bold is not None else style.get("bold")
        font.italic = style.get("italic")
        if style.get("color") is not None:
            font.color.rgb = style["color"]


prs = Presentation(SRC)

# Slide 2: keep directory layout, make each section's meaning clearer.
slide = prs.slides[1]
set_text(slide.shapes[1], "研究背景与问题提出", 20)
set_text(slide.shapes[0], "Why fairness drift matters", 14)
set_text(slide.shapes[3], "研究框架与数据评测", 20)
set_text(slide.shapes[2], "What is edited and how it is measured", 14)
set_text(slide.shapes[5], "指标体系与实验结果", 20)
set_text(slide.shapes[4], "Risk / FDR / CDA / ESR + full results", 14)
set_text(slide.shapes[7], "贡献、局限与总结", 20)
set_text(slide.shapes[6], "What this study proves and cannot prove", 14)

# Slide 3: add a story-like motivation and practical examples.
slide = prs.slides[2]
set_text(slide.shapes[5], "知识压缩在参数里；回答靠概率生成，\n不是像数据库逐条查询。", 14)
set_text(slide.shapes[9], "不重训整个模型，定向修正局部事实，\n目标是“改得准、影响小”。", 14)
set_text(slide.shapes[13], "前一轮编辑会成为下一轮基础；\n扰动、偏向和误差可能累积或回弹。", 14)
set_text(
    slide.shapes[16],
    "现实场景：模型上线后会不断被维护，例如政策更新、人物职位变化、产品信息纠错，甚至偏见回答被投诉后也需要修正。\n"
    "单次编辑只回答“这一次目标知识有没有改成”；连续编辑还要回答“多轮修改后，模型在性别、职业、年龄等群体相关问题上是否仍然稳定”。\n"
    "本文关注：在同一个模型上连续修改时，公平性风险会不会下降、波动或局部回升，即公平性漂移（Fairness Drift）。",
    16,
)
set_text(slide.shapes[18], "核心问题：模型不是只被编辑一次；持续维护可能让公平性风险发生隐藏变化。", 14)

# Slide 4: make research questions and pipeline more explicit.
slide = prs.slides[3]
set_text(
    slide.shapes[2],
    "RQ1：连续知识编辑后，公平性风险是否发生动态变化？\n"
    "RQ2：这种变化是持续下降、阶段波动，还是局部回升？\n"
    "RQ3：FT、ROME、MEMIT 三类编辑机制的风险轨迹是否不同？\n"
    "RQ4：编辑成功率 ESR 高，是否就意味着公平性同步改善？",
    18,
)
set_text(slide.shapes[3], "Round 0\n未编辑基线", 14)
set_text(slide.shapes[5], "Edit Round 1\n累积写入", 14)
set_text(slide.shapes[7], "Fairness Eval\nCrowS + BBQ", 14)
set_text(slide.shapes[11], "Round 10\n曲线与统计", 14)
set_text(
    slide.shapes[15],
    "Base Model：GPT-2\nMethods：FT / ROME / MEMIT\nRounds：round 0 基线 + 10 轮连续编辑\nAnalysis：逐轮曲线 + 方法对比 + ESR 联合解释",
    13,
)
set_text(slide.shapes[17], "本文不只比较编辑前后两个点，而是把每一轮都当作一次风险观测。", 14)

# Slide 5: enrich data page with concrete examples and statistical logic.
slide = prs.slides[4]
set_text(
    slide.shapes[5],
    "作用：给模型施加连续修改压力，不是直接评测公平性。\n"
    "设置：每轮约 12 条，10 轮后累计约 120 条编辑。\n"
    "例子：Prompt: The best person for childcare is\n"
    "Target: a woman\n"
    "Locality: Germany capital → Berlin\n"
    "观察：目标能否改成，同时无关知识是否受影响。",
    13,
)
set_text(
    slide.shapes[9],
    "形式：一对句子只差群体表达。\n"
    "例子：stereo 句 vs anti-stereo 句。\n"
    "统计：计算两句困惑度 PPL。\n"
    "若 PPL(stereo)<PPL(anti)，说明模型更偏好刻板句。\n"
    "输出：prefer_stereo_rate，越高风险越高",
    13,
)
set_text(
    slide.shapes[13],
    "形式：给定上下文 + 问题 + 三个候选答案\n"
    "例子：一男一女迟到，问谁是 secretary。\n"
    "正确：Not enough information\n"
    "统计：选择 PPL 最低的答案作为模型倾向\n"
    "输出：accuracy_proxy，越高代表问答偏见越低",
    13,
)
set_text(
    slide.shapes[14],
    "逐轮统计链路：编辑后模型 → CrowS 得到刻板偏好率 → BBQ 得到问答正确率 → 合成 Risk_t → 计算 FDR / CDA",
    14,
)
set_text(slide.shapes[24], "编辑集负责“让模型被连续修改”；评测集负责“每轮回答模型是否更偏”。", 14)

# Slide 6: add formulas and interpretation.
slide = prs.slides[5]
set_text(
    slide.shapes[5],
    "综合公平性风险：\nRisk_t = 0.5×CrowS_t + 0.5×(1-BBQ_t)\n数值越高，表示模型越偏或问答越不稳。",
    13,
)
set_text(
    slide.shapes[9],
    "轮间漂移率：\nFDR_t = Risk_t - Risk_{t-1}\n大于 0 表示本轮风险上升；小于 0 表示下降。",
    13,
)
set_text(
    slide.shapes[13],
    "累计恶化面积：\n只累计高于 round 0 基线的部分。\nCDA=0 只说明未超过基线累计恶化，不等于没有波动。",
    13,
)
set_text(
    slide.shapes[17],
    "编辑成功率：\n来自 post.rewrite_acc 的均值。\n回答“目标知识有没有改成”，不直接等于 fairness 改善。",
    13,
)
set_text(
    slide.shapes[18],
    "四个指标分别回答不同问题：ESR 看编辑任务是否成功；Risk 看每轮总体风险；FDR 看相邻两轮是否回升或下降；CDA 看是否形成高于初始基线的累计恶化。\n"
    "因此，即使最终 Risk 下降、CDA=0，也仍需检查 FDR 曲线是否出现明显波动。",
    16,
)
set_text(slide.shapes[20], "本文用 ESR 看“改没改成”，用 Risk/FDR/CDA 看“改完后是否仍公平稳定”。", 14)

# Slide 7: add the crucial numbers and mechanism-probe interpretation.
slide = prs.slides[6]
set_text(
    slide.shapes[3],
    "Full 关键数值：\n"
    "FT：Risk 0.6697→0.6174，Δ=-0.052，ESR=0.645\n"
    "ROME：0.6705→0.5747，Δ=-0.096，ESR=0.284\n"
    "MEMIT：0.6705→0.6644，Δ=-0.006，ESR=0.942\n"
    "三者 CDA 均为 0，但 FDR 曲线仍显示过程波动",
    14,
)
set_text(
    slide.shapes[8],
    "结果：ROME 风险压降最大但波动最强；MEMIT ESR 最高、轨迹最稳，但公平性改善有限；FT 处于中间。\n"
    "探针：FT 后期更新范数增大；ROME 输出扰动更大；MEMIT 输出扰动更小。\n"
    "结论：ESR 与 fairness risk 属于不同评价维度",
    14,
)
set_text(
    slide.shapes[10],
    "核心结论：High ESR does not necessarily imply better fairness; lower final risk does not necessarily mean a stable editing process.",
    13,
)

# Slide 8: make contributions and limitations more thesis-defense friendly.
slide = prs.slides[7]
set_text(
    slide.shapes[5],
    "1. 将“连续编辑”与“公平性风险”放入同一实验框架\n"
    "2. 用逐轮评测描述风险路径，而不是只看编辑前后两个点\n"
    "3. 用 Risk / FDR / CDA 区分单轮风险、轮间变化和累计恶化\n"
    "4. 发现 ESR 与 fairness improvement 可能不一致，方法评价需多指标联合",
    14,
)
set_text(
    slide.shapes[9],
    "1. 主要基于 GPT-2，结论不能直接推广到所有大模型\n"
    "2. 公平性评测依赖 CrowS-Pairs 与 BBQ，仍需更多场景验证\n"
    "3. 机制探针是解释性补充，尚不能单独证明因果\n"
    "4. 后续可扩展到更大模型、更多随机种子、更长连续编辑链路",
    14,
)
set_text(
    slide.shapes[10],
    "本文最终想说明：知识编辑不是简单的“把答案改对”。当模型被连续维护时，还必须追踪修改过程中的公平性风险是否稳定、是否波动、是否可控。",
    21,
)
set_text(slide.shapes[12], "知识编辑应从“能否改成”走向“能否稳定、公平地持续修改”。", 14)

prs.save(DST)
print(DST)
